

import streamlit as st
from pptx import Presentation
from pptx.util import Pt
import io

st.set_page_config(page_title="History of AI — Presentation", layout="wide")


slides = [
    {
        "title": "Introduction",
        "content": (
            "Artificial Intelligence (AI) enables machines to perform tasks requiring human intelligence.\n"
            "This presentation explores its evolution from early foundations to modern Agentic AI."
        ),
    },
    {
        "title": "Early Foundations (1940s–1950s)",
        "content": (
            "- 1943: McCulloch & Pitts propose first mathematical neuron model (NEURAL NETWORK CONCEPTS).\n"
            "- 1950: Alan Turing posses the questions, Can machines think?.\n"
            "- 1956: Dartmouth Conference – official birth of AI and coining the term Artificial Intelligence."
        ),
    },
    {
        "title": "The First AI Winter (1970s–1980s)",
        "content": (
            "- Unrealistic expectations led to funding cuts.\n"
            "- Limited computing power and data availability."
        ),
    },
    {
        "title": "The Knowledge-based Era (1980s)",
        "content": (
            "- AI revival via rule-based Expert Systems — Symbolic Reasoning &  Expert System.\n"
            "- Symbolic Reasoning — Attempting to replicate human intelligence using logical rules and explicit programming.\n"
            "- Expert System — Programs designed to mimic the decision-making ability of a human expert in a specific domain."
            "- Limitations — These systems struggled with ambiguity, lacked scalability, and were difficult to maintain."
        ),
    },
    {
        "title": "The Deep Learning Revolution (2000s - Present)",
        "content": (
            "- A major paradigm shift driven by massive data availability, increased computing power (GPUs), and algorithmic breakthroughs.\n"
            "- Faster Processing — The rise of powerful GPUs enabled the training of much deeper and larger neural networks..\n"
            "- Big Data Availability — Vast digital datasets became available, allowing models to learn complex patterns organically."
            "- Deep Learning models learn representations directly from raw data, eliminating the need for hand-crafted features"
        ),
    }, 
    
    {
        "title": "Understanding Large Language Models (LLMs)",
        "content": (
            "LLMs are advanced neural networks specifically designed for processing and generating human-like text at scale.\n"
            "Key ideas:\n"
            "- The core breakthrough is the Transformer Architecture (2017), which uses an attention mechanism to weigh the importance of different words in a sequence.\n"
            "- LLMs treat text as a sequence of tokens (symbols) and use their vast training to predict the next token in the sequence.\n"
            "- They function as a statistical engine for generating highly coherent and contextually relevant responses.\n" 
        ),         
    },

    {
        "title": "Stage 1: Pre-training (Creating the BaseModel)",
        "content": (
            "- Massive Data Scale:- Models are trained on an immense quantity(e.g., up to 15 trillion tokens) and diversity of public web text (e.g., Common Crawl).\n"
            "- Tokenization:- Raw text is converted into a finite sequence of symbolic tokens (e.g., via Byte Pair Encoding), which are the model9s basic units.\n"
            "- Next-Token Prediction:- The Transformer network is trained to constantly predict the next token in any given sequence, optimizing its billions of parameters.\n"
            "Important Point- Base Model is not yet a helpful assistant, but a statistical simulator.\n"           
        ),
    },

    {

        "title": "Stage 2: Supervised Fine-Tuning (SFT)",
        "content": (
            "- Curated Conversations:- The training data is swapped from general internet text to high-quality, human-curated conversational examples.\n"
            "- Human Labelers:- Contractors follow detailed, hundreds-of-page instructions to produce responses— helpful, truthful, and harmless.\n"
            "- Developing Identity:— This process shifts the model9s statistical identity, aligning its output with the desired behavior of a reliable assistant.\n"
        ),
    },
    {
        "title": "The Rise of Agentic AI",
        "content": (
            "- Planning and Goal Setting:- Agents break down complex, multi-step goals into smaller, manageable tasks before executing.\n"
            "- Tool Utilization:— LLMs are integrated with external tools like web search, code interpreters, and APIs to extend their capabilities.\n"
            "- Reflection and SelfCorrection:- The model can review its own output or progress, identify errors, and iterate on its approach to achieve a better result."
        ),
    },
    {
        "title": "From Static Models to Dynamic Agent",
        "content": (
            "- STATIC MODELS:-\n"
            "Generates output based solely on the initial prompt and internal knowledge.\n"
            "Limited to text generation; cannot directly interact with the real world or live data.\n"
            "Output is a single, static response; THINKING process is internal.\n"
            "- Dynamic Agents:-\n"            
            "Generates steps, executes actions in an external environment,and uses feedback."
            "Uses external APIs, browser interfaces, and file systems to complete tasks."
            "Output is a sequence of actions; the THINKING process (CoT) is transparent and iterative."
        ),
    },
    {
        "title": "Key Takeaways: The Future of Intelligence",
        "content": (
            "- Foundational Paradigm Shift:-\n"
            "AI has moved from rule-based systems to data-driven, representation-learning models (Deep Learning).\n"
            "- The Transformer Engine:-"
            "The attention mechanism remains the crucial breakthrough enabling the scale and performance of modern LLMs."
            "- Alignment is Key:-\n"
            "Post-training steps like Supervised Fine-Tuning are essential to align powerful base models with human values and utility."
            "- Agentic Autonomy:-"
            "The next major phase involves creating systems that can plan, act, and self-correct, greatly increasing the practicalutility of AI"
        ),
    },   
]

# --- UI ---
st.title("📚 The History & Evolution of Artificial Intelligence")
st.caption("Use the sidebar to navigate slides   or export a PPTX.")

slide_index = st.sidebar.number_input("Slide", min_value=1, max_value=len(slides), value=1, step=1)
show_notes = st.sidebar.checkbox("Show slide notes", value=True)

current = slides[slide_index - 1]
st.header(current["title"])
st.write(current["content"].replace("\n", "  \n"))

# show all slides as an accordion
with st.expander("Full Slide List"):
    for i, s in enumerate(slides, start=1):
        st.markdown(f"**{i}. {s['title']}**")
        if show_notes:
            st.write(s['content'])



def build_pptx(slides):
    prs = Presentation()
    for s in slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = s['title']
        body = slide.shapes.placeholders[1].text_frame
        body.text = s['content']
        for p in body.paragraphs:
            for run in p.runs:
                run.font.size = Pt(18)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

if st.sidebar.button("Export PPTX"):
    pptx_data = build_pptx(slides)
    st.download_button("Download Presentation (.pptx)", data=pptx_data, file_name="History_of_AI_Presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

st.markdown("---")
st.info("This presentation is for educational purposes. It summarizes major developments and does not cover every AI research contribution.")


st.write("**Credits:** Compiled from Asharib Ali'pdf  and written by ashhad.")
